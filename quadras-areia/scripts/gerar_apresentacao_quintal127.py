#!/usr/bin/env python3
"""Gera apresentação profissional — Quintal 127 | Sistema de Locação de Quadras."""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Quintal127_Apresentacao_Sistema.pptx"
LOGO = ROOT / "frontend" / "public" / "logo-quintal127.png"

# Paleta Quintal 127
BLUE = RGBColor(0x0D, 0x47, 0xA1)
BLUE_DARK = RGBColor(0x04, 0x1D, 0x40)
LIME = RGBColor(0x9A, 0xCD, 0x32)
LIME_LIGHT = RGBColor(0xE4, 0xF9, 0xC5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xF5)


def set_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9.1), Inches(0.85))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = LIME


def bullets(slide, left, top, width, height, items, size=14, title=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE
        p.space_after = Pt(8)
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def card(slide, left, top, w, h, title, lines, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), h - Inches(0.25))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = accent
    for line in lines:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(4)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLUE_DARK)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(1.5), Inches(1.2), width=Inches(7))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Sistema de Reserva e Locação de Quadras"
    p.font.size = Pt(22)
    p.font.color.rgb = LIME
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Rua Laureano 127, Campeche — Florianópolis, SC"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "Apresentação Técnica do Projeto"
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(24)


def content_slide(prs, title, subtitle, items, two_col=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_bar(slide, title, subtitle)
    if two_col:
        bullets(slide, Inches(0.5), Inches(1.4), Inches(4.3), Inches(5.5), two_col[0], size=13)
        bullets(slide, Inches(5.0), Inches(1.4), Inches(4.3), Inches(5.5), two_col[1], size=13, title=two_col[2] if len(two_col) > 2 else None)
    else:
        bullets(slide, Inches(0.55), Inches(1.45), Inches(8.9), Inches(5.5), items, size=14)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    cover_slide(prs)

    content_slide(prs, "Contexto", "O desafio do negócio", [
        "Empresas de locação de quadras precisam organizar reservas, pagamentos e agenda",
        "Processos manuais geram conflitos de horário e perda de receita",
        "Clientes esperam reservar online, de forma rápida e pelo celular",
        "A gestão interna exige controle de quadras, bloqueios e notas fiscais",
    ])

    content_slide(prs, "A Solução", "Quintal 127 — plataforma web completa", [
        "Sistema full-stack para reserva de quadras de areia",
        "Dois perfis: Cliente (reservas) e Administrador (gestão)",
        "Interface responsiva para web e mobile",
        "Pagamento simulado integrado ao fluxo de reserva",
        "Controle de disponibilidade em tempo real",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_bar(slide, "Stack Tecnológica", "Ferramentas modernas e consolidadas no mercado")
    stacks = [
        ("Frontend", ["React + Vite", "Tailwind CSS", "Zustand", "Axios", "React Router"]),
        ("Backend", ["Node.js + Express", "Prisma ORM", "JWT + bcrypt", "SQLite / PostgreSQL"]),
        ("Qualidade", ["API REST organizada", "Validações com Zod", "Middleware de segurança", "Seed de dados"]),
    ]
    for i, (title, lines) in enumerate(stacks):
        card(slide, Inches(0.5 + i * 3.15), Inches(1.5), Inches(2.95), Inches(2.8), title, lines)

    content_slide(prs, "Arquitetura do Sistema", "Separação clara de responsabilidades", [
        "Frontend consome API REST via Axios com token JWT",
        "Backend em camadas: routes → controllers → services → Prisma",
        "Banco relacional com models: User, Court, Booking, Payment, Invoice, BlockedTime",
        "Autenticação stateless com persistência no localStorage",
        "Deploy preparado para ambiente com internet",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_bar(slide, "Perfis de Acesso", "Cliente e Administrador")
    card(slide, Inches(0.55), Inches(1.55), Inches(4.2), Inches(4.8), "CLIENTE", [
        "Cadastro e login",
        "Ver quadras e horários",
        "Reservar e cancelar",
        "Pagar e ver comprovante",
        "Acompanhar nota fiscal",
    ], LIME)
    card(slide, Inches(5.1), Inches(1.55), Inches(4.2), Inches(4.8), "ADMINISTRADOR", [
        "Dashboard com métricas",
        "CRUD de quadras",
        "Gestão de reservas e clientes",
        "Pagamentos e notas fiscais",
        "Bloqueio de horários",
    ], BLUE)

    content_slide(prs, "Funcionalidades — Cliente", "Experiência do usuário final", [
        "Cadastro com nome, e-mail, CPF, telefone, CEP e data de nascimento",
        "Listagem de quadras ativas com preço por hora",
        "Consulta de horários disponíveis por data",
        "Reserva com cálculo automático do valor total",
        "Minhas reservas com status de pagamento e NF",
        "Cancelamento com regra de 24h de antecedência",
    ])

    content_slide(prs, "Funcionalidades — Admin", "Painel de gestão interna", [
        "Dashboard: total de reservas, reservas do dia, faturamento estimado",
        "Criar, editar e excluir quadras",
        "Filtrar reservas por data, quadra, cliente e status",
        "Alterar status de reserva e pagamento",
        "Emitir e registrar nota fiscal",
        "Bloquear horários específicos por quadra",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_bar(slide, "Fluxo de Reserva", "Da escolha ao pagamento confirmado")
    steps = [
        ("1", "Cliente escolhe quadra e data"),
        ("2", "Sistema exibe slots livres (08h–22h)"),
        ("3", "Reserva criada como PENDENTE"),
        ("4", "Pagamento via PIX, cartão ou dinheiro"),
        ("5", "Status CONFIRMADA + NF registrada"),
    ]
    for i, (num, text) in enumerate(steps):
        y = Inches(1.5 + i * 1.05)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = LIME
        circle.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.08), Inches(0.55), Inches(0.4))
        tb.text_frame.paragraphs[0].text = num
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.size = Pt(16)
        tb.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        line_box = slide.shapes.add_textbox(Inches(1.45), y + Inches(0.1), Inches(7.5), Inches(0.5))
        line_box.text_frame.paragraphs[0].text = text
        line_box.text_frame.paragraphs[0].font.size = Pt(16)
        line_box.text_frame.paragraphs[0].font.color.rgb = DARK

    content_slide(prs, "Regras de Negócio", "Garantias do sistema", [
        "Impossível reservar horário já ocupado na mesma quadra",
        "Horários bloqueados pelo admin ficam indisponíveis",
        "Valor total = horas × preço/hora da quadra",
        "Pagamento confirmado atualiza reserva para CONFIRMADA",
        "Cliente só cancela a própria reserva (24h antes)",
        "Admin pode cancelar qualquer reserva",
    ])

    content_slide(prs, "Modelagem de Dados", "Principais entidades (Prisma)", [
        "User — clientes e administradores",
        "Court — quadras com preço por hora",
        "Booking — reservas com data, horário e status",
        "Payment — PIX, cartão ou dinheiro",
        "Invoice — emissão de nota fiscal",
        "BlockedTime — bloqueios de agenda por quadra",
    ])

    content_slide(prs, "Segurança", "Boas práticas implementadas", [
        "Senhas criptografadas com bcrypt",
        "Autenticação JWT em rotas protegidas",
        "Middleware admin para área restrita",
        "Senhas nunca retornam nas respostas da API",
        "Validação de entrada nos controllers",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_bar(slide, "Interface e Identidade Visual", "Design alinhado à marca Quintal 127")
    card(slide, Inches(0.55), Inches(1.5), Inches(4.2), Inches(2.5), "Visual", [
        "Cores: azul profundo + verde limão",
        "Logo transparente no header",
        "Layout responsivo mobile/desktop",
        "Cards, badges de status e dashboard",
    ])
    card(slide, Inches(5.1), Inches(1.5), Inches(4.2), Inches(2.5), "Páginas", [
        "Home, Login, Cadastro",
        "Quadras, Reserva, Pagamento",
        "Painel Cliente e Painel Admin",
        "Gestão completa pelo navegador",
    ])
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(2.5), Inches(4.3), width=Inches(5))

    content_slide(prs, "Benefícios", "Para o negócio e para o cliente", [
        "Redução de conflitos de agenda",
        "Autonomia do cliente para reservar 24/7",
        "Visão financeira para o administrador",
        "Base escalável para novas quadras e filiais",
        "Experiência profissional alinhada à marca",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLUE_DARK)
    title_bar(slide, "Demonstração", "Como testar o sistema")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.5))
    tf = box.text_frame
    lines = [
        ("Repositório:", "github.com/dourado99k/VIBING-PRO-MAX (pasta quadras-areia)"),
        ("Frontend:", "localhost:5173"),
        ("API:", "localhost:3001"),
        ("Admin:", "admin@quadras.com / admin123"),
        ("Cliente:", "cadastre-se pela tela de registro"),
    ]
    for i, (label, value) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}  {value}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE if i == 0 else RGBColor(0xDD, 0xEE, 0xFF)
        p.space_after = Pt(14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLUE_DARK)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(2.2), Inches(1.5), width=Inches(5.6))
    box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado!"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = LIME
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Quintal 127 — Arena & Choperia | Rua Laureano 127, Campeche"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    prs.save(str(OUTPUT))
    print(f"✅ Apresentação gerada: {OUTPUT}")
    print(f"   Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

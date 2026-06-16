#!/usr/bin/env python3
"""Gera apresentação revisada — Estudo de Caso Lipedema (Paciente 01)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = "/Users/dourado99k/Downloads/Lipedema_Estagio_Clinico_REVISADO.pptx"
OUTPUT_ORIGINAL = "/Users/dourado99k/Downloads/Lipedema_Estagio_Clinico.pptx"

# Paleta
TEAL = RGBColor(0x0D, 0x6E, 0x6E)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF7, 0xF7)
ACCENT = RGBColor(0xE8, 0x6C, 0x4A)


def set_slide_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xCC, 0xEE, 0xEE)


def add_body(slide, left, top, width, height, lines, font_size=14, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.level = line[1]
        else:
            p.text = line
            p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(4)
    return box


def add_bullet_box(slide, left, top, width, height, title, items, font_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(1.5)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.15))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = TEAL
    for item in items:
        bp = tf.add_paragraph()
        bp.text = f"• {item}"
        bp.font.size = Pt(font_size)
        bp.font.color.rgb = DARK
        bp.space_before = Pt(3)


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight connector
    conn.line.color.rgb = TEAL
    conn.line.width = Pt(2)


def slide_title(prs, title, subtitle, extra=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TEAL)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCC, 0xEE, 0xEE)
    p2.alignment = PP_ALIGN.CENTER
    if extra:
        p3 = tf.add_paragraph()
        p3.text = extra
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(20)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 — Capa
    slide_title(
        prs,
        "Estudo de Caso Clínico — Lipedema",
        "Estácio de Sá · Curso de Nutrição",
        "Disciplina: Estágio de Nutrição Clínica\n"
        "Acadêmico: Valter Luis Bittencourt Mocinho Junior · Mat. 202303028989 · 2026",
    )

    # 2 — Estrutura
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Estrutura do Trabalho")
    topics = [
        "01  Dados do Paciente",
        "02  Patologia (Lipedema)",
        "03  Etiologia, Fisiopatologia e Sinais/Sintomas",
        "04  Antropometria e Limitações da Avaliação",
        "05  Diagnóstico Nutricional",
        "06  Sinais e Sintomas Clínicos",
        "07  Funcionamento do TGI",
        "08  Exames Bioquímicos e Interação Fármaco × Nutriente",
        "09  Recordatório Alimentar 24h",
        "10  Cálculo da Prescrição Dietoterápica",
        "11  Prescrição Dietoterápica",
        "12  Cardápio Prescrito",
        "13  Adequação: Prescrição × Conduta",
        "14  Adequação: Prescrição × R24h",
        "15  Orientações e Conclusão",
    ]
    add_body(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.8), topics, font_size=15)

    # 3 — Dados do paciente
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Dados do Paciente")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.4), Inches(4.3), Inches(2.8),
        "Identificação",
        [
            "Paciente 01",
            "Sexo feminino",
            "Data de nascimento: 05/02/2001 (25 anos)",
            "Diagnóstico clínico: Lipedema",
            "Atendimento: Estágio de Nutrição Clínica",
        ],
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.4), Inches(4.5), Inches(2.8),
        "Queixa principal",
        [
            "Desejo de melhorar composição corporal",
            "Queixas relacionadas ao acúmulo de gordura em membros inferiores",
            "Interesse em padrão alimentar mais adequado à condição",
        ],
    )
    add_body(
        slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(1.5),
        [
            "Observação ética: dados identificadores substituídos por sigla (Paciente 01), "
            "conforme orientação acadêmica para proteção da privacidade.",
        ],
        font_size=12,
    )

    # 4 — Patologia (slide exclusivo + referência)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Patologia: Lipedema")
    add_body(
        slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(4.5),
        [
            "Definição",
            "O lipedema é uma doença crônica do tecido adiposo, caracterizada por distribuição "
            "desproporcional e simétrica de gordura, predominantemente em membros inferiores "
            "(e, em alguns casos, braços), com preservação relativa de mãos e pés.",
            "",
            "Características principais",
            "• Distribuição ginoide da adiposidade, frequentemente independente do IMC",
            "• Pode coexistir com obesidade, mas é entidade clínica distinta",
            "• Associa-se a alterações microvasculares, disfunção linfática e processo inflamatório crônico de baixo grau",
            "• Impacto funcional, estético e na qualidade de vida",
            "• Predominância no sexo feminino; fatores hormonais e genéticos envolvidos",
            "",
            "Referência bibliográfica",
            "HERBST, K. L. et al. Lipedema: a review of the literature. International Journal of "
            "Endocrinology and Metabolism, v. 19, n. 3, e113052, 2021. "
            "https://doi.org/10.5812/ijem.113052",
        ],
        font_size=13,
        bold_first=False,
    )
    for i, txt in enumerate(slide.shapes[-1].text_frame.paragraphs):
        if txt.text in ("Definição", "Características principais", "Referência bibliográfica"):
            txt.font.bold = True
            txt.font.color.rgb = TEAL

    # 5 — Etiologia ↔ Fisiopatologia ↔ Sinais/Sintomas (interações)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Interações: Etiologia, Fisiopatologia e Sinais/Sintomas")
    add_bullet_box(
        slide, Inches(0.4), Inches(1.3), Inches(2.9), Inches(5.5),
        "ETIOLOGIA",
        [
            "Predisposição genética",
            "Alterações hormonais (estrogênio)",
            "Gatilhos: puberdade, gestação, menopausa",
            "Microangiopatia inicial",
        ],
        font_size=11,
    )
    add_bullet_box(
        slide, Inches(3.55), Inches(1.3), Inches(2.9), Inches(5.5),
        "FISIOPATOLOGIA",
        [
            "Hipertrofia/hiperplasia de adipócitos",
            "Disfunção linfática → edema intersticial",
            "Hipóxia tecidual local",
            "Processo inflamatório crônico de baixo grau (não é sinal/sintoma isolado)",
            "Fibrose progressiva do tecido adiposo",
            "Alterações hemodinâmicas e capilares",
        ],
        font_size=11,
    )
    add_bullet_box(
        slide, Inches(6.7), Inches(1.3), Inches(2.9), Inches(5.5),
        "SINAIS E SINTOMAS",
        [
            "Aumento de volume em pernas/braços (simétrico)",
            "Desproporção tronco × membros",
            "Dor à palpação / hipersensibilidade",
            "Equimoses espontâneas",
            "Sinal do caule (+)",
            "Pés e mãos poupadas",
            "Dificuldade de perda de gordura localizada",
        ],
        font_size=11,
    )
    add_body(
        slide, Inches(0.4), Inches(6.2), Inches(9.2), Inches(1.0),
        [
            "Interações: etiologia genético-hormonal → disfunção adiposa/linfática → inflamação crônica e fibrose → "
            "sinais clínicos (volume, dor, equimoses). A inflamação é processo patológico identificado clinicamente "
            "e por exames (ex.: PCR-us), não classificada como sinal ou sintoma isolado.",
        ],
        font_size=11,
    )

    # 6 — Antropometria
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Dados Antropométricos")
    metrics = [
        ("58,6 kg", "Peso"),
        ("1,56 m", "Altura"),
        ("24,1 kg/m²", "IMC (eutrofia)"),
        ("40,2%", "Gordura corporal (bioimpedância)"),
        ("71 cm", "Circunferência da cintura"),
        ("96 cm", "Circunferência do quadril"),
        ("0,74", "Relação cintura/quadril"),
    ]
    for i, (val, lbl) in enumerate(metrics):
        col, row = i % 4, i // 4
        x = Inches(0.5 + col * 2.35)
        y = Inches(1.5 + row * 1.6)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.1), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = TEAL
        tb = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2.1), Inches(1.0))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # 7 — Limitação dobras cutâneas
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Limitações da Avaliação Antropométrica")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(2.5),
        "Dobras cutâneas — não realizadas",
        [
            "As dobras cutâneas não foram aferidas no atendimento.",
            "Motivo: roupas inadequadas para a avaliação (tecido espesso, modelagem apertada ou "
            "cobertura insuficiente para marcação e pinçamento correto dos pontos anatômicos).",
            "A inadequação vestimentar impede a padronização da técnica e compromete a validade do resultado.",
        ],
    )
    add_bullet_box(
        slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(2.3),
        "Métodos utilizados",
        [
            "Peso e estatura (IMC)",
            "Bioimpedância elétrica (% de gordura corporal)",
            "Circunferências (cintura e quadril)",
        ],
    )

    # 8 — Diagnóstico nutricional (com parâmetros)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Diagnóstico Nutricional")
    add_body(
        slide, Inches(0.5), Inches(1.3), Inches(9.0), Inches(5.8),
        [
            "Avaliação dos parâmetros antropométricos coletados",
            "",
            "IMC 24,1 kg/m² → eutrofia (adequado segundo classificação para adultos)",
            "% Gordura corporal 40,2% → elevada para sexo feminino (referência: adequado ~20–32%; "
            "alterado >32%) → excesso de adiposidade",
            "Circunferência da cintura 71 cm → adequada (<80 cm para mulheres)",
            "Circunferência do quadril 96 cm → dentro de faixa observada",
            "Relação cintura/quadril 0,74 → adequada (<0,85 para mulheres)",
            "",
            "Diagnóstico nutricional (NCP)",
            "Excesso de adiposidade corporal relacionado ao lipedema e ao padrão alimentar inadequado, "
            "evidenciado por percentual de gordura elevado (40,2%) apesar de IMC em eutrofia.",
            "",
            "Etiologia relacionada: ingestão habitual de ultraprocessados e baixa adequação de fibras/vegetais.",
            "Sinais/sintomas relacionados: acúmulo de gordura em abdômen, pernas e tríceps; dificuldade de "
            "modificação da composição corporal.",
        ],
        font_size=13,
    )

    # 9 — Sinais e sintomas (separado)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Sinais e Sintomas Clínicos")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5),
        "Sinais (objetivos)",
        [
            "Acúmulo de gordura em abdômen, pernas e tríceps",
            "Distribuição desproporcional da adiposidade",
            "Equimoses espontâneas em membros",
            "Sinal do caule positivo",
            "Pés e mãos sem aumento de volume",
        ],
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.3), Inches(4.5), Inches(5.5),
        "Sintomas (subjetivos)",
        [
            "Dor ou desconforto à palpação em membros inferiores",
            "Sensação de peso nas pernas",
            "Hipersensibilidade ao toque",
            "Dificuldade de redução de gordura localizada com dieta",
        ],
    )
    add_body(
        slide, Inches(0.5), Inches(6.5), Inches(9.0), Inches(0.7),
        [
            "Nota: inflamação crônica é processo fisiopatológico/diagnóstico de suporte (ex.: PCR-us), "
            "não sendo classificada como sinal ou sintoma isolado.",
        ],
        font_size=11,
    )

    # 10 — TGI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Funcionamento do Trato Gastrointestinal (TGI)")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.4), Inches(9.0), Inches(4.5),
        "Avaliação clínica",
        [
            "Frequência evacuatória: diária",
            "Escala de Bristol: tipo 3–4 (fezes normais a pastosas)",
            "Ausência de queixas gastrointestinais no momento da consulta",
            "Sem relato de náuseas, vômitos, dor abdominal ou alteração recente do hábito intestinal",
            "Sem uso de laxantes ou probióticos",
        ],
        font_size=14,
    )

    # 11 — Exames
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Exames Bioquímicos e Interação Fármaco × Nutriente")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5),
        "Exames solicitados",
        [
            "Hemograma completo",
            "Glicemia de jejum",
            "Hemoglobina glicada (HbA1c)",
            "Insulina",
            "Perfil lipídico",
            "TGO / TGP",
            "Ureia e creatinina",
            "PCR-us (ultrassensível)",
            "VHS",
        ],
        font_size=12,
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.3), Inches(4.5), Inches(5.5),
        "Interação fármaco × nutriente",
        [
            "Paciente 01 não utiliza medicamentos de uso contínuo.",
            "Sem interações medicamentosas a monitorar no plano alimentar.",
            "Exames pendentes de coleta para acompanhamento do estado inflamatório e metabólico.",
        ],
        font_size=13,
    )

    # 12 — R24h
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Recordatório Alimentar 24h (R24h)")
    meals = [
        ("☀️ Café da manhã / Desjejum", "YoPRO (iogurte proteico)"),
        ("🍽️ Almoço", "Frango, arroz, macarrão, lasanha, batata frita, salada e refrigerante"),
        ("🥪 Lanche", "Rap10 com ovo, queijo e presunto"),
        ("🌙 Jantar", "Miojo com ovo e queijo"),
    ]
    for i, (meal, content) in enumerate(meals):
        y = Inches(1.35 + i * 1.35)
        add_bullet_box(slide, Inches(0.5), y, Inches(9.0), Inches(1.2), meal, [content], font_size=13)

    # 13 — Cálculo da prescrição
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Cálculo da Prescrição Dietoterápica")
    add_body(
        slide, Inches(0.45), Inches(1.25), Inches(9.1), Inches(6.0),
        [
            "Dados utilizados: Paciente 01 — 25 anos, sexo feminino, 58,6 kg, 1,56 m",
            "",
            "1) Taxa Metabólica Basal (TMB) — Mifflin-St Jeor (mulheres):",
            "   TMB = (10 × 58,6) + (6.25 × 156) − (5 × 25) − 161 = 1.275 kcal/dia",
            "",
            "2) Gasto Energético Total (GET) — atividade leve (FA = 1,375):",
            "   GET = 1.275 × 1,375 = 1.753 kcal/dia",
            "",
            "3) Valor Energético Total (VET) prescrito:",
            "   VET = 1.163 kcal/dia (déficit energético de ~34% em relação ao GET, visando recomposição corporal)",
            "",
            "4) Proteínas prescritas:",
            "   71,0 g/dia = 1,21 g/kg (58,6 kg)",
            "",
            "5) Distribuição de macronutrientes prescrita (valores do plano alimentar):",
            "   PTN: 71,0 g (284 kcal — 24,4%)  |  LIP: 34,9 g (314 kcal — 27,0%)  |  CHO: 156,4 g (626 kcal — 53,8%)",
            "",
            "6) Fibras prescritas: 28,6 g/dia  |  Hidratação orientada: 1,8 L/dia",
        ],
        font_size=12,
    )

    # 14 — Prescrição (macros exatos)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Prescrição Dietoterápica")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(2.5),
        "Metas prescritas (valores reais do plano)",
        [
            "VET: 1.163 kcal/dia",
            "Proteínas: 71,0 g/dia",
            "Lipídeos: 34,9 g/dia",
            "Carboidratos: 156,4 g/dia",
            "Fibras: 28,6 g/dia",
            "Hidratação: 1,8 L/dia",
        ],
        font_size=13,
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.3), Inches(4.5), Inches(2.5),
        "Conduta nutricional",
        [
            "Padrão alimentar com foco anti-inflamatório",
            "Ômega-3 (salmão/tilápia) e frutas/vegetais",
            "Redução de ultraprocessados",
            "Distribuição proteica ao longo do dia",
            "Fibras adequadas para saciedade e TGI",
        ],
        font_size=13,
    )
    add_body(
        slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(2.5),
        [
            "Micronutrientes prescritos (totais diários): Cálcio 313,7 mg · Magnésio 262,2 mg · Ferro 9,0 mg · "
            "Vit. C 64,5 mg · Vit. D 6,4 mcg · Sódio 1.073,7 mg · Potássio 2.313,0 mg",
            "Perfil lipídico da dieta: AG monoinsaturados 9,6 g · poliinsaturados 9,6 g · saturados 5,9 g · "
            "trans 0,2 g · colesterol 124,2 mg",
        ],
        font_size=11,
    )

    # 15 — Cardápio café
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Cardápio Prescrito — Café da Manhã e Lanche")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5),
        "Café da manhã (opção principal)",
        [
            "Abacate — 3 colheres de sopa (135 g)",
            "Pão de forma tradicional — 2 fatias (50 g)",
            "",
            "Substituições:",
            "Avocado — 2 col. sopa cheias (90 g) ou 1 porção grupo 10",
            "Torrada integral — 4 unidades (40 g) ou 2 porções grupo 12",
        ],
        font_size=12,
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.3), Inches(4.5), Inches(5.5),
        "Lanche (opção principal)",
        [
            "Torrada integral — 4 unidades (40 g)",
            "1 porção do grupo Açúcares (lista de substituição)",
            "",
            "Substituições:",
            "Pão de forma tradicional — 2 fatias (50 g) ou 2 porções grupo 12",
        ],
        font_size=12,
    )

    # 16 — Cardápio almoço e jantar
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Cardápio Prescrito — Almoço e Jantar")
    add_bullet_box(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.8),
        "Almoço (opção principal)",
        [
            "Couve refogada — 2 col. sopa cheias picada (40 g)",
            "Cebola cozida — 4 col. sopa cheias picada (40 g)",
            "Brócolis cozido — 2 ramos médios (120 g)",
            "Cenoura baby — 6 unidades médias (60 g)",
            "Filé de salmão com pele grelhado — 0,5 filé pequeno (60 g)",
            "Arroz branco cozido — 2 col. sopa cheias (50 g)",
            "Feijão preto cozido — 2 conchas pequenas cheias (130 g)",
            "Banana prata — 1 unidade média (65 g)",
        ],
        font_size=11,
    )
    add_bullet_box(
        slide, Inches(5.0), Inches(1.3), Inches(4.5), Inches(5.8),
        "Jantar (opção principal)",
        [
            "Pepino cru — 2 unidades médias (200 g)",
            "Tomate salada — 6 fatias grandes (180 g)",
            "Filé de tilápia cozido — 2 filés médios (120 g)",
            "Arroz branco cozido — 1,5 col. sopa cheia (38 g)",
            "Grão-de-bico cozido — 1 col. arroz cheia (45 g)",
        ],
        font_size=11,
    )

    # 17 — Adequação prescrição × conduta
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Adequação: Prescrição × Conduta Nutricional")
    rows = [
        ("Objetivo clínico", "Prescrição", "Adequação"),
        ("Reduzir inflamação", "Ômega-3 (salmão, tilápia); fibras 28,6 g; vegetais variados", "Adequado"),
        ("Recomposição corporal", "VET 1.163 kcal; PTN 71 g (1,21 g/kg)", "Adequado"),
        ("Saciedade / TGI", "Fibras 28,6 g; vegetais em todas as refeições", "Adequado"),
        ("Reduzir ultraprocessados", "Alimentos in natura e minimamente processados", "Adequado"),
        ("Distribuição proteica", "PTN em café, almoço, lanche e jantar", "Adequado"),
        ("Hidratação", "Orientação 1,8 L/dia", "Adequado"),
    ]
    tbl = slide.shapes.add_table(len(rows), 3, Inches(0.4), Inches(1.35), Inches(9.2), Inches(4.5)).table
    col_w = [Inches(2.2), Inches(4.5), Inches(2.5)]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if r > 0 else 12)
                p.font.bold = r == 0
                p.font.color.rgb = TEAL if r == 0 else DARK

    # 18 — Adequação prescrição × R24h
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Adequação: Prescrição × Recordatório 24h")
    rows2 = [
        ("Aspecto", "R24h (habitual)", "Prescrição", "Adequação"),
        ("Energia", "Ultraprocessados e porções amplas (lasanha, fritas, refrigerante)", "1.163 kcal/dia", "R24h acima do prescrito"),
        ("Proteínas", "Presentes, mas com fontes processadas", "71,0 g/dia", "Qualidade proteica inferior no R24h"),
        ("Fibras", "Baixa (poucos vegetais)", "28,6 g/dia", "R24h inadequado"),
        ("Ultraprocessados", "Rap10, miojo, refrigerante, lasanha, batata frita", "Ausentes no cardápio", "R24h inadequado"),
        ("Vegetais", "Apenas salada no almoço", "Vegetais em almoço e jantar", "R24h inadequado"),
        ("Fontes de gordura", "Frituras e embutidos", "Salmão, abacate, preparações grelhadas/cozidas", "R24h inadequado"),
        ("Refrigerante", "Presente no almoço", "Não prescrito", "R24h inadequado"),
    ]
    tbl2 = slide.shapes.add_table(len(rows2), 4, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.0)).table
    widths = [Inches(1.5), Inches(2.5), Inches(2.8), Inches(2.6)]
    for i, w in enumerate(widths):
        tbl2.columns[i].width = w
    for r, row_data in enumerate(rows2):
        for c, cell_text in enumerate(row_data):
            cell = tbl2.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10 if r > 0 else 11)
                p.font.bold = r == 0
                p.font.color.rgb = TEAL if r == 0 else DARK

    # 19 — Orientações
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Orientações Nutricionais")
    orient = [
        ("Hidratação", "Mínimo de 1,8 L de água por dia, distribuídos ao longo do dia."),
        ("Frutas e vegetais", "Consumo diário conforme cardápio prescrito; priorizar variedade de cores."),
        ("Ultraprocessados", "Reduzir refrigerantes, embutidos, fast food, miojo e produtos industrializados."),
        ("Preparo dos alimentos", "Priorizar grelhados, cozidos e refogados; evitar frituras."),
        ("Atividade física", "Manter prática regular para apoio à composição corporal e circulação."),
    ]
    for i, (title, desc) in enumerate(orient):
        y = Inches(1.35 + i * 1.1)
        add_bullet_box(slide, Inches(0.5), y, Inches(9.0), Inches(0.95), title, [desc], font_size=12)

    # 20 — Conclusão
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TEAL)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusão"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    items = [
        "A terapia nutricional é fundamental no manejo do lipedema.",
        "A prescrição de 1.163 kcal/dia com 71 g de proteínas atende aos objetivos de recomposição corporal.",
        "O padrão anti-inflamatório prescrito contrasta com o R24h, que apresenta excesso de ultraprocessados.",
        "O acompanhamento dos exames bioquímicos e da composição corporal é essencial.",
        "A intervenção visa melhorar composição corporal, modular inflamação e promover qualidade de vida.",
    ]
    for item in items:
        bp = tf.add_paragraph()
        bp.text = f"• {item}"
        bp.font.size = Pt(15)
        bp.font.color.rgb = WHITE
        bp.space_before = Pt(8)
    foot = tf.add_paragraph()
    foot.text = "\nValter Luis Bittencourt Mocinho Junior · Mat. 202303028989 · Estácio de Sá · 2026"
    foot.font.size = Pt(12)
    foot.font.color.rgb = RGBColor(0xCC, 0xEE, 0xEE)
    foot.alignment = PP_ALIGN.CENTER
    foot.space_before = Pt(24)

    prs.save(OUTPUT)
    prs.save(OUTPUT_ORIGINAL)
    print(f"Salvo: {OUTPUT}")
    print(f"Atualizado: {OUTPUT_ORIGINAL}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
